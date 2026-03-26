#!/usr/bin/env python3
"""Export SVG slides to PPTX. Default: PNG embed (universal). --svg: SVG embed (Office 365+)."""

import argparse
import re
import sys
import zipfile
from io import BytesIO
from pathlib import Path

import cairosvg
from lxml import etree
from pptx import Presentation
from pptx.util import Inches


SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)
PNG_WIDTH = 2560
PNG_HEIGHT = 1440


def discover_slides(slides_dir: Path) -> list[Path]:
    """Find slide-{nn}.svg files sorted by number."""
    pattern = re.compile(r"^slide-(\d+)\.svg$")
    matched = []
    for f in slides_dir.iterdir():
        m = pattern.match(f.name)
        if m:
            matched.append((int(m.group(1)), f))
    matched.sort(key=lambda t: t[0])
    if not matched:
        print(f"Error: no slide-*.svg files found in {slides_dir}", file=sys.stderr)
        sys.exit(1)
    return [f for _, f in matched]


def parse_speaker_notes(notes_path: Path) -> dict[int, str]:
    """Parse speaker notes markdown. Format: '## Slide N\\n\\nNotes text\\n\\n'."""
    notes = {}
    if not notes_path.exists():
        return notes
    text = notes_path.read_text(encoding="utf-8")
    blocks = re.split(r"^## Slide\s+(\d+)\s*$", text, flags=re.MULTILINE)
    # blocks: ['preamble', '1', 'notes1', '2', 'notes2', ...]
    for i in range(1, len(blocks) - 1, 2):
        slide_num = int(blocks[i])
        content = blocks[i + 1].strip()
        if content:
            notes[slide_num] = content
    return notes


def svg_to_png(svg_path: Path) -> bytes:
    """Render SVG to PNG at 2x retina resolution."""
    svg_data = svg_path.read_bytes()
    return cairosvg.svg2png(
        bytestring=svg_data,
        output_width=PNG_WIDTH,
        output_height=PNG_HEIGHT,
    )


def export_png_mode(
    slide_files: list[Path],
    output_path: Path,
    notes: dict[int, str],
) -> None:
    """Export slides as PNG-embedded PPTX (universal compatibility)."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank

    for idx, svg_path in enumerate(slide_files, start=1):
        png_data = svg_to_png(svg_path)
        slide = prs.slides.add_slide(blank_layout)
        img_stream = BytesIO(png_data)
        slide.shapes.add_picture(
            img_stream, left=0, top=0, width=SLIDE_WIDTH, height=SLIDE_HEIGHT
        )
        if idx in notes:
            slide.notes_slide.notes_text_frame.text = notes[idx]

    prs.save(str(output_path))


def export_svg_mode(
    slide_files: list[Path],
    output_path: Path,
    notes: dict[int, str],
) -> None:
    """Export slides with SVG embedded in OOXML (Office 365+ only).

    Strategy: build a standard PPTX first, then inject SVG into the ZIP package.
    Each slide gets an SVG image relationship and the picture element references it.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # Build slides with PNG fallback images first
    slide_png_map: list[tuple[int, bytes, Path]] = []
    for idx, svg_path in enumerate(slide_files, start=1):
        png_data = svg_to_png(svg_path)
        slide = prs.slides.add_slide(blank_layout)
        img_stream = BytesIO(png_data)
        slide.shapes.add_picture(
            img_stream, left=0, top=0, width=SLIDE_WIDTH, height=SLIDE_HEIGHT
        )
        if idx in notes:
            slide.notes_slide.notes_text_frame.text = notes[idx]
        slide_png_map.append((idx, png_data, svg_path))

    # Save to buffer, then reopen as ZIP to inject SVGs
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    out_buf = BytesIO()
    nsmap = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "asvg": "http://schemas.microsoft.com/office/drawing/2016/SVG/main",
        "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    }

    with zipfile.ZipFile(buf, "r") as zin, zipfile.ZipFile(
        out_buf, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)

            # Inject SVG relationship into slide .rels files
            slide_match = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels$", item.filename)
            if slide_match:
                slide_num = int(slide_match.group(1))
                if 1 <= slide_num <= len(slide_files):
                    svg_filename = f"../media/slide-{slide_num:02d}.svg"
                    rels_tree = etree.fromstring(data)
                    # Find max rId
                    max_rid = 0
                    for rel in rels_tree:
                        rid = rel.get("Id", "")
                        m = re.match(r"rId(\d+)", rid)
                        if m:
                            max_rid = max(max_rid, int(m.group(1)))
                    new_rid = f"rId{max_rid + 1}"
                    svg_rel = etree.SubElement(rels_tree, "Relationship")
                    svg_rel.set("Id", new_rid)
                    svg_rel.set(
                        "Type",
                        "http://schemas.microsoft.com/office/2016/09/relationships/svgImage",
                    )
                    svg_rel.set("Target", svg_filename)
                    data = etree.tostring(rels_tree, xml_declaration=True, encoding="UTF-8")

                    # Also patch the slide XML to add asvg:svgBlip
                    slide_xml_path = f"ppt/slides/slide{slide_num}.xml"
                    slide_data = zin.read(slide_xml_path)
                    slide_tree = etree.fromstring(slide_data)
                    # Find blipFill/blip and add svgBlip extension
                    for blip in slide_tree.iter(f"{{{nsmap['a']}}}blip"):
                        ext_lst = blip.find(f"{{{nsmap['a']}}}extLst")
                        if ext_lst is None:
                            ext_lst = etree.SubElement(blip, f"{{{nsmap['a']}}}extLst")
                        ext = etree.SubElement(ext_lst, f"{{{nsmap['a']}}}ext")
                        ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
                        svg_blip = etree.SubElement(ext, f"{{{nsmap['asvg']}}}svgBlip")
                        svg_blip.set(f"{{{nsmap['r']}}}embed", new_rid)
                    patched_slide_data = etree.tostring(
                        slide_tree, xml_declaration=True, encoding="UTF-8"
                    )
                    zout.writestr(slide_xml_path, patched_slide_data)

            # Skip slide XML files we already wrote patched versions of
            if re.match(r"ppt/slides/slide\d+\.xml$", item.filename):
                slide_num_check = int(re.search(r"slide(\d+)", item.filename).group(1))
                if 1 <= slide_num_check <= len(slide_files):
                    continue  # already written above

            zout.writestr(item, data)

        # Add SVG files to media/
        for idx, svg_path in enumerate(slide_files, start=1):
            svg_data = svg_path.read_bytes()
            zout.writestr(f"ppt/media/slide-{idx:02d}.svg", svg_data)

        # Update [Content_Types].xml to include SVG content type
        content_types_data = zin.read("[Content_Types].xml")
        ct_tree = etree.fromstring(content_types_data)
        ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
        # Check if .svg extension already registered
        has_svg = any(
            el.get("Extension") == "svg"
            for el in ct_tree.iter(f"{{{ct_ns}}}Default")
        )
        if not has_svg:
            svg_default = etree.SubElement(ct_tree, f"{{{ct_ns}}}Default")
            svg_default.set("Extension", "svg")
            svg_default.set("ContentType", "image/svg+xml")
        zout.writestr(
            "[Content_Types].xml",
            etree.tostring(ct_tree, xml_declaration=True, encoding="UTF-8"),
        )

    out_buf.seek(0)
    output_path.write_bytes(out_buf.read())


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Export SVG slides to PPTX presentation."
    )
    parser.add_argument("slides_dir", type=Path, help="Directory containing slide-{nn}.svg files")
    parser.add_argument("output", type=Path, help="Output .pptx file path")
    parser.add_argument(
        "--svg",
        action="store_true",
        help="Embed SVG directly in OOXML (Office 365+ only)",
    )
    parser.add_argument(
        "--speaker-notes",
        type=Path,
        default=None,
        help="Markdown file with speaker notes (## Slide N format)",
    )
    args = parser.parse_args()

    if not args.slides_dir.is_dir():
        print(f"Error: {args.slides_dir} is not a directory", file=sys.stderr)
        sys.exit(1)

    slide_files = discover_slides(args.slides_dir)
    notes = parse_speaker_notes(args.speaker_notes) if args.speaker_notes else {}

    mode = "SVG-embed" if args.svg else "PNG-embed"
    print(f"Exporting {len(slide_files)} slides ({mode} mode)...")

    args.output.parent.mkdir(parents=True, exist_ok=True)

    if args.svg:
        export_svg_mode(slide_files, args.output, notes)
    else:
        export_png_mode(slide_files, args.output, notes)

    size_mb = args.output.stat().st_size / (1024 * 1024)
    print(f"Done: {len(slide_files)} slides exported to {args.output} ({size_mb:.1f} MB)")


if __name__ == "__main__":
    main()
